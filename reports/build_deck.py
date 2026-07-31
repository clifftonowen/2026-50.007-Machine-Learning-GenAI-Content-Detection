"""Generate reports/presentation.pptx, the Task 4 presentation.

Same numbers as the report, from reports/facts.py, so the two deliverables cannot
disagree. Slides carry the claim; the chosen-versus-tried reasoning lives in the speaker
notes so nothing on screen is a wall of text.

    python reports/build_deck.py
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import facts  # noqa: E402

OUT = facts.PROJECT_ROOT / "reports" / "presentation.pptx"

W, H = Inches(13.333), Inches(7.5)          # 16:9
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GOOD = RGBColor(0x2C, 0xA0, 0x2C)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]

    def _text(self, slide, left, top, width, height, text, *, size=18,
              bold=False, color=INK, align=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = par.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            if align is not None:
                par.alignment = align
        return box

    def slide(self, title, *, subtitle=None, notes=""):
        s = self.prs.slides.add_slide(self.blank)
        self._text(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9),
                   title, size=30, bold=True)
        if subtitle:
            self._text(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.6),
                       subtitle, size=15, color=MUTED)
        if notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    def bullets(self, slide, items, *, top=2.0, size=17, left=0.7, width=12.0):
        height = min(4.4, H / 914400 - top - 0.3)      # never run past the slide edge
        assert height > 0.5, f"no room for bullets at top={top}"
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                       Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            par.space_after = Pt(10)
            run = par.add_run()
            run.text = item if item.startswith(" ") else "- " + item
            run.font.size = Pt(size)
            run.font.color.rgb = INK
        return box

    def picture(self, slide, name, *, top=1.95, height=4.9, left=None):
        path = facts.figure(name)
        pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top),
                                       height=Inches(height))
        pic.left = Emu(int((W - pic.width) / 2)) if left is None else Inches(left)
        return pic

    def stat(self, slide, value, label, *, left, top=2.4, color=ACCENT, size=54):
        self._text(slide, Inches(left), Inches(top), Inches(4.0), Inches(1.0),
                   value, size=size, bold=True, color=color)
        self._text(slide, Inches(left), Inches(top + 1.0), Inches(4.0), Inches(1.0),
                   label, size=13, color=MUTED)

    def save(self):
        self.prs.save(OUT)
        return OUT


def build():
    facts.check()
    d = Deck()

    # 1 title
    s = d.slide("Detecting Machine-Generated Text",
                subtitle="50.007 Machine Learning  |  COLING 2026 GenAI Content Detection",
                notes="One sentence: we got most of our gain from what the model sees "
                      "and from how we turn its scores into labels, not from the model.")
    d._text(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(0.8),
            f"Public leaderboard Macro F1  {facts.BEST_KAGGLE}", size=40, bold=True,
            color=ACCENT)
    d._text(s, Inches(0.6), Inches(3.6), Inches(12.1), Inches(1.6),
            "LightGBM on 40,385 features built from raw text,\n"
            "with the two test populations thresholded separately.", size=19,
            color=MUTED)
    d._text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6),
            ", ".join(facts.MEMBERS), size=14, color=MUTED)

    # 2 the task
    s = d.slide("The task, and the trap inside it",
                notes="The shift is deliberate. High training F1 and low test F1 is the "
                      "designed behaviour, so a large gap is not evidence of a bug. What "
                      "it does mean is that same-domain validation will mislead you.")
    d.bullets(s, [
        "Binary: is this document machine-generated? Scored on Macro F1.",
        "20,000 labelled training documents, 6,999 to predict. 62.52% machine.",
        "Train is HC3 + M4GT + MAGE. Test is CUDRT + IELTS + NLPeer + PeerSum + MixSet.",
        "Zero corpus overlap. The shift is by design, and the brief says so.",
        f"Public leaderboard noise floor: {facts.NOISE_FLOOR}. Everything smaller is a tie.",
    ], top=2.1, size=19)

    # 3 the journey
    s = d.slide("The journey, in seven milestone submissions",
                subtitle="Representation, then calibration. The model itself contributed "
                         "almost nothing.",
                notes="Point at the two steep steps. Everything flat is a round we spent "
                      "on tuning or ensembling. We report those as nulls.")
    d.picture(s, "kaggle_journey.png", top=1.85, height=5.0)

    # 4 the data
    s = d.slide("What the data looked like",
                notes="Class imbalance forced stratification and class weighting "
                      "everywhere. Length correlates with the label, which later became "
                      "our validation protocol.")
    d.picture(s, "class_balance.png", top=2.0, height=4.3, left=0.9)
    d.picture(s, "text_length_distribution.png", top=2.0, height=4.3, left=6.9)

    # 5 baselines
    s = d.slide("Where we started: twelve models on the supplied features",
                subtitle="LightGBM led, and kept leading once we compared fairly.",
                notes="Chose LightGBM. Tried XGBoost, HistGB, three linear models, "
                      "naive Bayes, KNN, RandomForest, AdaBoost. Drawback of a single "
                      "model is no diversity. Acceptable because we later measured what "
                      "ensembling was worth and it was inside the noise floor.")
    d.picture(s, "baseline_models_cv_f1.png", top=2.0, height=4.7)

    # 6 wrong turn 1
    s = d.slide("Wrong turn 1: we compared models along an axis we had not fixed",
                notes="This is the difficulties question. We wrote down a conclusion, "
                      "designed a test that could overturn it, and it did. Keep that "
                      "habit visible.")
    d.bullets(s, [
        "For three notebooks we believed local validation was anti-correlated with Kaggle.",
        "LightGBM had the best CV and holdout of anything, and the worst leaderboard "
        "score. So we wrote it off.",
        "The cause: each model was submitted at whatever class balance its own default "
        "threshold produced.",
        "At a matched share, LightGBM beat ElasticNet by 0.0189, matching its local lead.",
        "Local validation had never been broken. We had been asking two questions as one.",
    ], top=2.1, size=18)

    # 7 the null round
    s = d.slide("Rounds 3 and 4: ensembling, and an honest null",
                subtitle="Four combiner families, 48 configurations, +0.0017 on Kaggle.",
                notes="A fifth of the noise floor. The diagnosis was that we had been "
                      "optimising the model while holding the representation fixed.")
    d.picture(s, "ensemble_combiner_leaderboard.png", top=2.0, height=4.7)

    # 8 the diagnosis
    s = d.slide("The diagnosis: we were tuning the wrong thing",
                notes="The supplied features are lemmatised with stop words removed, so "
                      "they are pure content. Content vocabulary is exactly what changes "
                      "between corpora. Function words, punctuation and layout survive.")
    d.bullets(s, [
        "The supplied features are top-5,000 TF-IDF over lemmas, stop words removed.",
        "That is a pure content signal, and content is what differs between corpora.",
        "Raw text was fully available and essentially untouched: one length summary in "
        "seven notebooks.",
        "So we rebuilt the representation from the text: function words, punctuation, "
        "casing, layout, length, diversity, and character and word n-grams.",
    ], top=2.1, size=19)

    # 9 representation result
    s = d.slide("Changing the representation: +0.036 on the leaderboard",
                subtitle="Every dense block alone scores below the supplied control. "
                         "Together they beat it by 0.115.",
                notes="Chose all raw-text blocks. Tried supplied alone 0.7303, style "
                      "only 0.8458, char n-grams alone 0.8188, supplied plus everything "
                      "0.8795 which ties raw text alone with 5,000 more columns. "
                      "Drawback: 40,385 hand-built features. Acceptable because the "
                      "supplied ones delete the signal that transfers.")
    d.picture(s, "representation_comparison.png", top=2.0, height=4.7)

    # 10 falsified prediction
    s = d.slide("Wrong turn 2: we predicted the formatting features were artifacts",
                notes="Machine text in training carries markdown bold at ten times the "
                      "human rate, and the peer-review test rows carry it at five times "
                      "the machine rate. Sounded like a domain marker. We built the test "
                      "so it could fail, and it did.")
    d.bullets(s, [
        "Hypothesis: layout and punctuation features are dataset fingerprints and will "
        "not survive the corpus change.",
        "Test: ship a stripped-down feature set with only function words and character "
        "n-grams.",
        "Result: it lost 0.04944, six times the noise floor, in the opposite direction.",
        "We would rather report a falsified prediction than a roadmap where every idea "
        "worked.",
    ], top=2.1, size=19)

    # 11 validation
    s = d.slide("Which local number should we trust?",
                subtitle="Standard CV overstates the leaderboard by ~0.10. Grouped CV "
                         "gets within ~0.03.",
                notes="Chose leave-one-length-band-out. Tried standard 5-fold, k-means "
                      "domain clusters which would not reproduce across machines, and a "
                      "random grouping as a control. The control is the point: it "
                      "reproduced standard CV, so the gap comes from structure and not "
                      "from training on fewer rows.")
    d.picture(s, "local_vs_kaggle.png", top=2.1, height=4.4)

    # 12 the share finding
    s = d.slide("The biggest single finding: threshold the two populations separately",
                subtitle="A global cut moves both groups the same way. They need "
                         "opposite corrections.",
                notes="We swept global share twice, on two models, and it was flat both "
                      "times. Both measurements were correct and both conclusions were "
                      "wrong. A +0.019 effect and a -0.019 effect were cancelling.")
    d.picture(s, "share_surface.png", top=2.1, height=4.4)

    # 13 the general lesson
    s = d.slide("Why the global sweep looked flat",
                notes="This is the idea that generalises furthest beyond this project. "
                      "Say it slowly.")
    d._text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.2),
            "When one global threshold is applied to a population that is a mixture,\n"
            "the optimum for the mixture can be flat while the component optima are\n"
            "far apart and moving in opposite directions.",
            size=26, bold=True, color=ACCENT)
    d.bullets(s, [
        "UUID rows wanted a higher predicted machine share. Numeric rows wanted a lower "
        "one.",
        "Correcting them separately: +0.022 on the leaderboard, no model change at all.",
        "Confirmed independently on a second, unrelated model to within 0.001.",
    ], top=4.7, size=17)

    # 14 what did not work
    s = d.slide("What did not work, and we report it",
                subtitle="Four feature ideas, all inside the selection bar.",
                notes="Extended diversity actually made transfer worse while improving "
                      "same-domain fit, which is the exact memorisation signature. "
                      "Selecting on standard CV alone would have shipped it.")
    d.bullets(s, [f"{n}: {v:+.4f} on held-out-domain Macro F1"
                  for n, v in facts.EXPANSION_NULLS]
              + ["K-means domain clusters: passed the stability gate on one machine at "
                 "ARI 0.9703 and failed on another at 0.5991.",
                 "Ensembling on the supplied features: +0.0017, a fifth of the noise "
                 "floor."],
              top=2.2, size=18)

    # 15 attribution
    s = d.slide("Where the gain actually came from",
                notes="Two thirds representation, one third calibration, essentially "
                      "nothing from the model. That is not the split we expected when we "
                      "started.")
    d.stat(s, f"{facts.SUPPLIED_LGBM_KAGGLE}", "supplied TF-IDF, global threshold",
           left=0.8, color=MUTED)
    d.stat(s, "+0.044", "raw-text representation", left=4.9)
    d.stat(s, "+0.022", "per-group thresholding", left=9.0, color=GOOD)
    d._text(s, Inches(0.8), Inches(4.9), Inches(11.8), Inches(1.4),
            f"Final: {facts.BEST_KAGGLE}    "
            f"Local: {facts.BEST_CV_STANDARD} standard CV, "
            f"{facts.BEST_CV_GROUPED} held-out domain", size=20, bold=True)
    d._text(s, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.0),
            "Hyperparameter tuning contributed nothing so far: the shipped model runs on "
            "LightGBM defaults.", size=15, color=MUTED)

    # 16 lessons and next
    s = d.slide("What we learned, and what is still open",
                notes="Close on the methodology, not the score. The paired-comparison "
                      "point and the mixture point are the two we would defend hardest.")
    d.bullets(s, [
        "Read the dataset paper first. It produced both changes that moved our score.",
        "When folds differ by design, compare paired differences, not means.",
        "A leaderboard has a precision. Compute it once, then hold every claim to it.",
        "Averaging can hide two large effects that cancel.",
        " ",
        "Open: hyperparameter search on this representation, a second model family, and "
        "final pick selection.",
        "Final picks pair the best public point with a more central one, because our "
        "share coordinates are fitted to a noisy leaderboard.",
    ], top=2.1, size=17)

    return d.save()


if __name__ == "__main__":
    out = build()
    print(f"wrote {out}  ({out.stat().st_size / 1000:.0f} kB)")
