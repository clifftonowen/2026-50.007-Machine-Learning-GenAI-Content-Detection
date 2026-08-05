"""Generate reports/report_draft.pptx, the Task 4 presentation told as a methodology.

Eighteen slides for a fifteen-minute slot: title, the task, the nine phases in order, and
three closing slides. Every slide carries speaker notes, and the notes hold the pros, the
cons and the tradeoff sentence so nothing on screen is a wall of text.

Numbers come from reports/facts.py, the same source as report_draft.docx.

    python reports/build_deck_draft.py
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import facts  # noqa: E402

OUT = facts.PROJECT_ROOT / "reports" / "report_draft.pptx"

W, H = Inches(13.333), Inches(7.5)          # 16:9
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GOOD = RGBColor(0x2C, 0xA0, 0x2C)
WARN = RGBColor(0xC0, 0x39, 0x2B)


def _text_height(items, size, width_in, *, space_after_pt=9):
    """Estimate the rendered height of a text block, in inches.

    Calibri averages a little under half the point size per character. We use 0.53
    deliberately, which is wider than the true average, so the estimate errs toward
    predicting too many lines. An assertion guarding against overflow is only useful if it
    is pessimistic; an optimistic one passes on exactly the slides that break.
    """
    chars_per_line = max(1, int(width_in * 72 / (0.53 * size)))
    lines = sum(max(1, -(-len(t) // chars_per_line)) for t in items)
    return (lines * 1.2 * size + len(items) * space_after_pt) / 72


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]

    def _text(self, slide, left, top, width, height, text, *, size=18,
              bold=False, color=INK, align=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = par.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            if align is not None:
                par.alignment = align
        return box

    def slide(self, title, *, phase=None, subtitle=None, notes=""):
        s = self.prs.slides.add_slide(self.blank)
        top = Inches(0.35)
        if phase:
            self._text(s, Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.4),
                       phase.upper(), size=11, bold=True, color=ACCENT)
            top = Inches(0.62)
        self._text(s, Inches(0.6), top, Inches(12.1), Inches(0.9), title,
                   size=28, bold=True)
        if subtitle:
            self._text(s, Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.6),
                       subtitle, size=14.5, color=MUTED)
        assert notes, f"every slide needs speaker notes: {title}"
        s.notes_slide.notes_text_frame.text = notes
        return s

    def bullets(self, slide, items, *, top=2.1, size=17, left=0.7, width=12.0,
                color=INK, limit=7.3):
        """Add a bullet block and return the inch position of its estimated bottom.

        PowerPoint does not clip an overflowing textbox, it just draws over whatever is
        below, and there is no renderer here to eyeball the result. So the height is
        estimated from the text and asserted against `limit`: a slide that would spill
        fails the build instead of shipping quietly broken.
        """
        used = _text_height(items, size, width)
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                       Inches(used))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            par.space_after = Pt(9)
            run = par.add_run()
            run.text = item if item.startswith(" ") else "- " + item
            run.font.size = Pt(size)
            run.font.color.rgb = color
        bottom = top + used
        assert bottom <= limit, (
            f"bullets from {top:.2f}in run to {bottom:.2f}in, past the {limit}in limit")
        return bottom

    def options(self, slide, rows, *, top=2.0, size=14):
        """Alternatives as a compact three-column block: option, pro, con."""
        left, widths = 0.7, [3.9, 4.0, 4.0]
        for j, head in enumerate(["Considered", "For", "Against"]):
            self._text(slide, Inches(left + sum(widths[:j])), Inches(top),
                       Inches(widths[j]), Inches(0.35), head, size=12, bold=True,
                       color=MUTED)
        y = top + 0.45
        for option, pro, con in rows:
            lines = max(len(option) // 32, len(pro) // 34, len(con) // 34) + 1
            block = 0.34 * lines + 0.16
            for j, (txt, colour) in enumerate([(option, INK), (pro, GOOD),
                                               (con, WARN)]):
                self._text(slide, Inches(left + sum(widths[:j])), Inches(y),
                           Inches(widths[j] - 0.2), Inches(block), txt, size=size,
                           color=colour)
            y += block
        assert y < 7.3, f"options block overflows the slide, ends at {y:.2f}in"
        return y

    def chose(self, slide, text, *, top):
        """The closing line of a phase slide, placed under whatever came before it."""
        used = _text_height(["We chose: " + text], 15, 12.0)
        assert top + used <= 7.4, (
            f"the chosen-line at {top:.2f}in would run to {top + used:.2f}in")
        box = self._text(slide, Inches(0.7), Inches(top), Inches(12.0), Inches(used),
                         "We chose: " + text, size=15, bold=True, color=ACCENT)
        return box

    def picture(self, slide, name, *, top=1.95, height=4.9, left=None):
        path = facts.figure(name)
        pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top),
                                       height=Inches(height))
        pic.left = Emu(int((W - pic.width) / 2)) if left is None else Inches(left)
        assert Emu(int(top * 914400)) + pic.height <= H + Inches(0.1), \
            f"{name} runs off the bottom of the slide"
        return pic

    def stat(self, slide, value, label, *, left, top=2.6, color=ACCENT, size=52):
        self._text(slide, Inches(left), Inches(top), Inches(4.0), Inches(1.0),
                   value, size=size, bold=True, color=color)
        self._text(slide, Inches(left), Inches(top + 1.0), Inches(4.0), Inches(1.0),
                   label, size=13, color=MUTED)

    def save(self):
        self.prs.save(OUT)
        return OUT


def build():
    facts.check()
    d = Deck()

    # 1 -------------------------------------------------------------- title
    s = d.slide(
        "Detecting Machine-Generated Text",
        subtitle="50.007 Machine Learning  |  COLING 2026 GenAI Content Detection",
        notes="Open with the method, not the score. We ran nine phases in order and each "
              "one is presented the same way: what we considered, what we chose, and the "
              "tradeoff we accepted. Two of the nine returned nothing, and we will say so.")
    d._text(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(0.8),
            f"Public leaderboard Macro F1  {facts.BEST_KAGGLE}", size=38, bold=True,
            color=ACCENT)
    d._text(s, Inches(0.6), Inches(3.5), Inches(12.1), Inches(1.6),
            f"LightGBM on {facts.BEST_N_FEATURES:,} features built from raw text,\n"
            "with the two test populations thresholded separately.", size=19, color=MUTED)
    d._text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6),
            ", ".join(facts.MEMBERS), size=14, color=MUTED)

    # 2 --------------------------------------------------------------- task
    s = d.slide("The task, and the trap inside it",
                notes="The shift is deliberate, so a large train-to-test gap is the "
                      "designed behaviour rather than a bug. What it does mean is that "
                      "same-domain validation will mislead you, and that is the thread "
                      "running through the whole deck.")
    d.bullets(s, [
        "Binary: is this document machine-generated? Scored on Macro F1.",
        "20,000 labelled training documents, 6,999 to predict. 62.52% machine.",
        "Train is HC3 + M4GT + MAGE. Test is CUDRT + IELTS + NLPeer + PeerSum + MixSet.",
        "Zero corpus overlap. The shift is by design and the brief says so.",
        f"Public leaderboard noise floor: {facts.NOISE_FLOOR}. Anything smaller is a tie.",
    ], top=1.9, size=19)

    # 3 ------------------------------------------------------------- method
    s = d.slide("Our method, in the order we ran it",
                subtitle="Four foundation phases, two that returned nothing, three that "
                         "moved the score.",
                notes="This slide is the spine of the talk. Point out that phases 5 and 6 "
                      "are grey on purpose: tuning and ensembling both came back inside "
                      "the noise floor, and that null is what redirected us to phase 8.")
    d.picture(s, "methodology_flow.png", top=1.9, height=5.3)

    # 4 ------------------------------------------------------------ phase 1
    s = d.slide("Measure the data before modelling it", phase="Phase 1  |  notebook 01",
                notes="Considered: start modelling immediately, versus full EDA first, "
                      "versus both in parallel. For EDA-first: stratification, class "
                      "weighting and metric choice all become measured decisions. "
                      "Against: a whole notebook with no score at the end. Acceptable "
                      "because three of four measurements became constraints we never "
                      "revisited, and the fourth was worth more than any model.")
    bottom = d.bullets(s, [
        "62.52% machine, 37.48% human. Predict-majority scores 62.5% accuracy and "
        "0.385 Macro F1.",
        "So: Macro F1 not accuracy, stratify every split, class_weight balanced "
        "everywhere.",
        "The matrix is 98.64% zero. About 68 of 5,000 columns active per document.",
        "The test file holds two populations: 1,999 UUID ids and 5,000 numeric ids.",
        "The numeric rows are peer reviews. Longer, heavier markdown, 15.8% contain "
        "\"reviewer\" against 0.2%.",
    ], top=1.9, size=17)
    d.chose(s, "a full exploratory pass before training anything.", top=bottom + 0.15)

    # 5 ------------------------------------------------------------ phase 2
    s = d.slide("Dimensionality reduction, and why we stopped using it",
                phase="Phase 2  |  notebooks 01 and 03  |  Tasks 1 and 2",
                subtitle="PCA where the task requires it. Nowhere else.",
                notes="This slide carries Task 2 and the component analysis its top band "
                      "asks for. Considered: PCA everywhere, TruncatedSVD, feature "
                      "selection, or nothing. Against PCA downstream: it is lossy rather "
                      "than efficient here. Acceptable because boosted trees take sparse "
                      "input directly and select features as they split.")
    bottom = d.bullets(s, [
        "Task 2 fixes PCA at 2,000 / 1,000 / 500 / 100 components with KNN, "
        "n_neighbors = 2.",
        "  Variance retained: " + ", ".join(f"{n:,} keeps {v:.1%}"
                                            for n, v in facts.PCA_VARIANCE) + ".",
        f"  No canonical threshold is reached within 2,000 components. The knee sits at "
        f"{facts.PCA_ELBOW[0]} components and keeps only {facts.PCA_ELBOW[1]:.1%}.",
        "Why: 98.64% zeros and ~68 active columns per document leaves little shared "
        "linear structure to absorb.",
        "That also predicts KNN finishing tenth of twelve: distances are weak in a sparse "
        "space, and worse after discarding 84% of the variance.",
    ], top=2.0, size=16)
    d.chose(s, "no dimensionality reduction for any Task 3 model.", top=bottom + 0.15)

    # 6 ------------------------------------------------------------ phase 3
    s = d.slide("Lock the validation protocol", phase="Phase 3  |  notebook 01",
                subtitle="The class balance from phase 1 decides this, not preference.",
                notes="At 62.5 against 37.5, an unstratified fold drifts from the "
                      "population by several points, and Macro F1 is sensitive to exactly "
                      "that drift. Tradeoff: it trains and tests inside the same corpora, "
                      "so it overstates the leaderboard by about 0.10 and can say nothing "
                      "about the test set's class balance. Acceptable because we needed a "
                      "ranking instrument, not a forecasting one, and it ranked correctly "
                      "every time we checked.")
    bottom = d.options(s, [
        ("Single random split", "Cheapest, one number per model",
         "Too noisy to separate adjacent candidates"),
        ("Plain k-fold", "Uses all the data",
         "Fold balance drifts, and Macro F1 is sensitive to that"),
        ("Stratified 5-fold", "Every fold matches the population; gives a spread too",
         "Five fits per candidate; same-corpus only"),
        ("Stratified 10-fold", "Lower variance still",
         "Twice the compute for a decimal we could not use"),
    ], top=2.15)
    d.chose(s, "stratified 5-fold, seed 42, Macro F1, plus a 4,000-row holdout saved to "
               "disk so all five of us evaluate identical rows.", top=bottom + 0.15)

    # 7 ------------------------------------------------------------ phase 4
    s = d.slide("Baselines before a single hyperparameter moved",
                phase="Phase 4  |  notebook 04",
                subtitle="Twelve families at library defaults. Only then did we tune.",
                notes="The spread is the finding. Twelve families span 0.14 and the top "
                      "four sit within 0.015, which told us in advance that model choice "
                      "was worth less than we hoped. That is the evidence we used later "
                      "to stop tuning and stop ensembling rather than grinding at them.")
    d.picture(s, "baseline_models_cv_f1.png", top=2.0, height=4.5)

    # 8 ------------------------------------------------------------ phase 5
    s = d.slide("Tuning worked locally, and then went badly",
                phase="Phase 5  |  notebooks 05 and 06",
                notes="Two-stage search: coarse randomised, then a focused grid. Every "
                      "local diagnostic said the model was healthy. The holdout agreed "
                      "with CV to 0.0006, so we were not overfitting the search. Then the "
                      "leaderboard came back 0.079 lower.")
    d.bullets(s, [
        f"LightGBM cross-validation {facts.SUPPLIED_LGBM_CV} to {facts.TUNED_LGBM_CV} "
        f"after the search.",
        f"Untouched holdout confirmed it at {facts.TUNED_LGBM_HOLDOUT}. Gap of 0.0006, so "
        "no search overfitting.",
        "Threshold tuned on the holdout added a further 0.003 locally.",
        f"Kaggle: {facts.FIRST_KAGGLE_DEFAULT_THR} at the default cut, "
        f"{facts.FIRST_KAGGLE_TUNED_THR} with the tuned one.",
        "A gap of roughly 0.079, and our carefully tuned threshold moved the leaderboard "
        "by under 0.01.",
    ], top=2.0, size=18)

    # 9 ------------------------------------------- phase 5, how we reacted
    s = d.slide("How we reacted to a 0.079 gap", phase="Phase 5  |  the decision",
                notes="This is the most reusable decision in the project. Computing the "
                      "noise floor cost one submission and paid for itself three times: "
                      "it stopped the ensembling phase, it stopped the share search "
                      "inside a flat region, and it made us disbelieve two of our own "
                      "gains.")
    bottom = d.options(s, [
        ("Trust local, keep tuning", "Local validation was clean and reproducible",
         "Local and Kaggle disagreed by 0.079"),
        ("Assume a bug, audit", "Cheap to check, severe if real",
         "The brief predicts this gap by design"),
        ("Tune against the leaderboard", "It is the graded surface",
         "~3,570 rows, rationed submissions: this is how leaderboard overfitting starts"),
        ("Keep both, measure the leaderboard",
         "Makes the two instruments comparable",
         "Spends submissions on measurement, not candidates"),
    ], top=2.0)
    d.chose(s, f"keep local validation for ranking, audit for leakage once, and compute "
               f"the leaderboard noise floor at {facts.NOISE_FLOOR}.", top=bottom + 0.15)

    # 10 ----------------------------------------------------------- phase 6
    s = d.slide("Ensembling, and an honest null", phase="Phase 6  |  notebooks 07 and 10",
                subtitle="Six members, four combiner families, forty-eight "
                         "configurations.",
                notes=f"Out-of-fold AUC improved 0.0075. Kaggle moved "
                      f"{facts.ROUND4_BEST_KAGGLE - facts.SUPPLIED_LGBM_KAGGLE:+.5f}, "
                      f"under a fifth of the noise floor. Without the noise floor we "
                      f"would have read that as a small win and added more members. With "
                      f"it, the correct reading was that two consecutive phases of model "
                      f"work had returned nothing, so the constraint was somewhere else.")
    d.picture(s, "ensemble_combiner_leaderboard.png", top=2.1, height=4.3)

    # 11 ----------------------------------------------------------- phase 7
    s = d.slide("Threshold calibration, and a confound we had been carrying",
                phase="Phase 7  |  notebooks 08, 09 and 11",
                notes="Considered: keep 0.5, tune on the holdout, sweep predicted share, "
                      "or set share from the paper. The holdout inherits the training "
                      "balance, so it optimises for the wrong distribution: measured at "
                      "0.66502. Tradeoff of sweeping share is that it fits the public "
                      "rows. Acceptable because every alternative is simply wrong, and we "
                      "manage the risk against the noise floor.")
    d.bullets(s, [
        "A 0.5 cut silently assumes the test balance equals the training balance of "
        "62.5%. The brief says it does not.",
        "Tuning the threshold on our holdout moved it the wrong way: 0.66502.",
        "No local data can estimate the test class balance. Dev and holdout are both "
        "carved from training.",
        "So we reparameterised: sort by score, label a fixed share machine. Comparable "
        "across models, and sweepable.",
        "That exposed the confound: every model had been submitted at whatever share its "
        "own default threshold produced.",
        f"At a matched share, LightGBM beat ElasticNet by 0.0189, matching its local lead. "
        "Local validation had never been broken.",
        f"Moving the tuned LightGBM off the default cut: "
        f"{facts.FIRST_KAGGLE_DEFAULT_THR} to {facts.SUPPLIED_LGBM_KAGGLE}. The largest "
        "single gain of the project, from one line of code.",
    ], top=1.9, size=16)

    # 12 -------------------------------------------------- phase 8, the why
    s = d.slide("The pivot: we had never questioned the input",
                phase="Phase 8  |  reading the paper properly",
                notes="This is the turning point of the talk. Three phases had changed "
                      "the model or the decision rule while holding the input fixed. What "
                      "changed our minds was the shared-task paper's data section, which "
                      "we had skimmed the first time.")
    d.bullets(s, [
        "The paper: train is HC3 + M4GT + MAGE, test is CUDRT + IELTS + NLPeer + PeerSum "
        "+ MixSet, no overlap.",
        "It also gives the balances: 62.6% machine in train against 53.1% in test. Both "
        "matched what we had measured.",
        "The supplied features are top-5,000 TF-IDF over lemmas with stop words removed.",
        "That is a pure content signal, and content vocabulary is exactly what changes "
        "when the corpus changes.",
        "Function words, punctuation, casing and layout survive a corpus change. The "
        "course preprocessing deletes precisely those.",
        "The raw text had been sitting in train.csv, untouched for seven notebooks.",
    ], top=1.9, size=16.5)

    # 13 ------------------------------------------------ phase 8, the result
    s = d.slide("Rebuilding the representation: +0.044 on the leaderboard",
                phase="Phase 8  |  notebooks 12 to 15",
                subtitle="Every dense block alone scores below the supplied control. "
                         "Together they beat it by 0.115.",
                notes="Considered: keep the supplied features, add a few style features, "
                      "rebuild n-grams only, or rebuild fully. Tradeoff: 40,385 "
                      "hand-built columns against 5,000 supplied, far more code, and the "
                      "brief awards no extra marks for own preprocessing. Acceptable "
                      "because it was worth +0.044 on the graded surface, which is more "
                      "than everything else we did to the model combined.")
    d.picture(s, "representation_comparison.png", top=2.1, height=4.3)

    # 14 --------------------------------------------- phase 8, falsification
    s = d.slide("We predicted the formatting features were artifacts. They were not.",
                phase="Phase 8  |  a test built to be able to fail",
                notes="Machine text in training carries markdown bold at ten times the "
                      "human rate, and the peer-review test rows carry it at five times "
                      "the machine rate. That sounded exactly like a dataset fingerprint. "
                      "We built the stripped-down set specifically so the prediction "
                      "could be falsified, and it was.")
    d.bullets(s, [
        "Hypothesis: layout and punctuation are dataset fingerprints and will not survive "
        "the corpus change.",
        "Test: ship a stripped set of function words and character n-grams only.",
        "Result: it lost 0.04944, six times the noise floor, in the opposite direction.",
        "Second falsified prediction: extending the diversity block made transfer worse "
        "while improving same-domain fit, which is the memorisation signature.",
        "We would rather report two falsified predictions than a roadmap where every idea "
        "worked.",
    ], top=2.0, size=17)

    # 15 ----------------------------------------------------------- phase 9
    s = d.slide("Calibration revisited, one share per population",
                phase="Phase 9  |  notebook 16",
                subtitle=f"uuid {facts.BEST_SHARES['uuid']}, numeric "
                         f"{facts.BEST_SHARES['numeric']}. No change to the model at all.",
                notes="Considered: one global share, separate shares per group, a model "
                      "per group, or domain adaptation. A model per group is impossible: "
                      "1,999 unlabelled UUID test rows. Tradeoff: both coordinates are "
                      "fitted to the public leaderboard over seven submissions. "
                      "Acceptable because the first correction gained 0.01764, a "
                      "deliberate wrong-direction probe lost 0.00926, and the same "
                      "correction reproduced on a different model.")
    d.picture(s, "share_surface.png", top=2.2, height=4.2)

    # 16 ------------------------------------------------------ the big idea
    s = d.slide("Why the global sweep looked flat", phase="The idea that generalises",
                notes="Say this one slowly. We measured the flat curve twice, correctly "
                      "both times, and drew the wrong conclusion on both occasions. The "
                      "only way to see it was to stop averaging.")
    d._text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.2),
            "When one global threshold is applied to a population that is a mixture,\n"
            "the optimum for the mixture can be flat while the component optima are\n"
            "far apart and moving in opposite directions.",
            size=25, bold=True, color=ACCENT)
    d.bullets(s, [
        "UUID rows wanted a higher predicted machine share. Numeric rows wanted a lower "
        "one.",
        "Roughly +0.019 and roughly -0.019, cancelling in the average.",
        "Correcting them separately: +0.022 on the leaderboard, with no model change.",
        "Confirmed independently on a second, unrelated model to within 0.001.",
    ], top=4.5, size=16)

    # 17 ------------------------------------------------------- attribution
    s = d.slide("Where the gain actually came from",
                notes="The single cheapest change we made was also the biggest: one line "
                      "moving off the default 0.5 cut, worth +0.078, nearly twice what "
                      "rebuilding the whole representation returned. If a marker asks "
                      "one question, expect it to be this one.")
    d.stat(s, "+0.078", "global share, off the default cut (phase 7)", left=0.8,
           color=GOOD, size=48)
    d.stat(s, "+0.044", "raw-text representation (phase 8)", left=4.9, size=48)
    d.stat(s, "+0.022", "per-group thresholding (phase 9)", left=9.0, color=GOOD,
           size=48)
    d._text(s, Inches(0.8), Inches(4.85), Inches(11.8), Inches(1.4),
            f"{facts.FIRST_KAGGLE_DEFAULT_THR} to {facts.BEST_KAGGLE}. "
            "About 70% from how scores become labels, 30% from what the model sees.",
            size=19, bold=True)
    d._text(s, Inches(0.8), Inches(5.6), Inches(11.8), Inches(1.4),
            "Nothing measurable from the model itself. Ensembling returned +0.0017 and "
            "the shipped model runs on LightGBM defaults.\n"
            "This is not the split any of us expected when we started.",
            size=15, color=MUTED)

    # 18 ---------------------------------------------- limitations and next
    s = d.slide("Limitations, and what is still open",
                notes="Close on the limitations rather than the score. The two we would "
                      "defend hardest are that share cannot be validated locally at all, "
                      "and that our grouped protocol ranks but does not forecast. Finish "
                      "on the final-pick rule.")
    d.bullets(s, [
        "Predicted share cannot be validated locally. Every share decision rests on "
        "leaderboard feedback, and we managed that risk rather than solving it.",
        "The grouped protocol is trustworthy to about 0.03 in level. It ranks candidates; "
        "it does not forecast a score.",
        "A ~0.075 standard-CV to leaderboard gap remains. The brief says to expect it.",
        "Open: hyperparameter search on this representation, designed and split five "
        "ways, not yet run. A second model family, specified but unmeasured.",
        " ",
        "Final picks pair the best public point with a more central one, because our "
        "share coordinates are fitted to a leaderboard with a known noise floor.",
    ], top=1.9, size=16.5)

    return d.save()


if __name__ == "__main__":
    out = build()
    print(f"wrote {out}  ({out.stat().st_size / 1000:.0f} kB)")
